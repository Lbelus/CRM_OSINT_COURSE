import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import sys

# Helper function to set up slide background and header/footer
def setup_slide(slide, title_text, slide_number):
    # Add footer with page number
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.5)
    footer = slide.shapes.add_textbox(left, top, width, height)
    text_frame = footer.text_frame
    p = text_frame.add_paragraph()
    p.text = f"{title_text} | Page {slide_number}"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT

    # Optional: Add a logo or image in the footer
    if os.path.exists("logo.png"):
        slide.shapes.add_picture("logo.png", Inches(0.5), Inches(6.8), width=Inches(1))

# Add a title slide
def add_title_slide(presentation, title, subtitle):
    slide_layout = presentation.slide_layouts[0]  # 'Title Slide' layout
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    # Set title properties
    title_placeholder.text = title
    title_paragraph = title_placeholder.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(40)
    title_paragraph.font.bold = True

    # Set subtitle properties
    subtitle_placeholder.text = subtitle
    subtitle_paragraph = subtitle_placeholder.text_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(24)

# Add slides with detailed content
def add_detailed_slides(presentation, title, slide_contents):
    for content in slide_contents:
        slide_layout = presentation.slide_layouts[1]  # 'Title and Content' layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set up slide background and footer
        setup_slide(slide, title, len(presentation.slides))

        # Set title
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        title_paragraph = title_placeholder.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(22)
        title_paragraph.font.bold = True

        # Set content
        content_placeholder = slide.placeholders[1]
        content_placeholder.text_frame.clear()  # Clear existing content

        # Add content paragraphs with indentation
        for line in content.split('\n'):
            p = content_placeholder.text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(16)
            if line.startswith('  ') or line.startswith('*'):
                p.level = 1  # Indent bullet points to a nested level
            else:
                p.level = 0  # Main bullet points

# Add section header slides
def add_section_header_slide(presentation, title, description=None):
    slide_layout = presentation.slide_layouts[2]  # 'Section Header' layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set up slide background and footer
    setup_slide(slide, title, len(presentation.slides))

    # Set title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_paragraph = title_placeholder.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True

    # Use the second placeholder (if available) for the description
    if description and len(slide.placeholders) > 1:
        description_placeholder = slide.placeholders[1]
        description_placeholder.text = description
        description_paragraph = description_placeholder.text_frame.paragraphs[0]
        description_paragraph.font.size = Pt(18)
        description_paragraph.alignment = PP_ALIGN.CENTER

# Add image slide
def add_image_slide(presentation, image):
    slide_layout = presentation.slide_layouts[6]  # 'Blank' layout for image slides
    slide = presentation.slides.add_slide(slide_layout)

    # Set up slide background and footer
    setup_slide(slide, "Image", len(presentation.slides))

    # Add image to the slide
    if os.path.exists(image):
        left = Inches(1.5)
        top = Inches(1)
        slide.shapes.add_picture(image, left, top, width=Inches(7))

# Add end slide
def add_end_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[5]  # 'Title Slide' layout for end slide
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_paragraph = title_placeholder.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True

    # Set content
    text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

# Main function to read JSON and create presentation
def create_presentation_from_json(json_file):
    with open(json_file, 'r') as file:
        data = json.load(file)

    presentation = Presentation()
    added_images = set()  # To keep track of added images

    # Add title slide
    add_title_slide(presentation, data['title_slide']['title'], data['title_slide']['subtitle'])

    # Add content and intermediary slides
    for slide_data in data['slides']:
        if slide_data['type'] == 'content':
            add_detailed_slides(presentation, slide_data['title'], slide_data['contents'])
            # Add image slide if an image is provided and hasn't been added yet
            image = slide_data.get('image')
            if image and os.path.exists(image) and image not in added_images:
                add_image_slide(presentation, image)
                added_images.add(image)
        elif slide_data['type'] == 'intermediary':
            add_section_header_slide(presentation, slide_data['title'], slide_data.get('description'))
            # Add image slide if an image is provided and hasn't been added yet
            image = slide_data.get('image')
            if image and os.path.exists(image) and image not in added_images:
                add_image_slide(presentation, image)
                added_images.add(image)

    # Add end slide
    add_end_slide(presentation, data['end_slide']['title'], data['end_slide']['content'])

    # Save the presentation
    presentation.save('Generated_Presentation.pptx')
    print("Presentation created successfully: Generated_Presentation.pptx")

# Run the program
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python create_presentation.py <path_to_json_file>")
    else:
        create_presentation_from_json(sys.argv[1])

